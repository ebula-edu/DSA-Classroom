import os
import math
from PySide6.QtCore import Qt, QSize, QPointF, QRectF
from PySide6.QtGui import QPixmap, QImage, QFont, QPainter, QColor, QPen, QBrush, QKeyEvent, QPainterPath
from PySide6.QtWidgets import (QWidget, QVBoxLayout, QHBoxLayout, QLabel, 
                             QPushButton, QScrollArea, QTextBrowser, 
                             QStackedWidget, QSplitter)
from gui.icons import VectorIconProvider

# Try importing dependencies
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from docx import Document
except ImportError:
    Document = None


class AnnotationOverlay(QWidget):
    """
    A transparent overlay that sits on top of document/slide images.
    Allows drawing, highlighting, and pointing without modifying underlying files.
    """
    def __init__(self, parent=None):
        super().__init__(parent)
        # Capture mouse events but draw transparently
        self.setAttribute(Qt.WA_NoSystemBackground, True)
        self.setAttribute(Qt.WA_TranslucentBackground, True)
        
        self.paths = []  # list of {"path": QPainterPath, "pen": QPen}
        self.current_path = None
        self.mode = "none"  # none (click-through), pen, highlight, pointer, eraser
        self.laser_pos = None
        self.set_mode("none")

    def set_mode(self, mode):
        self.mode = mode
        if mode == "none":
            self.setAttribute(Qt.WA_TransparentForMouseEvents, True) # Click-through
            self.setCursor(Qt.ArrowCursor)
        elif mode == "pointer":
            self.setAttribute(Qt.WA_TransparentForMouseEvents, False)
            self.setCursor(Qt.BlankCursor) # Hide cursor to show laser dot
        else:
            self.setAttribute(Qt.WA_TransparentForMouseEvents, False)
            self.setCursor(Qt.CrossCursor)
        self.update()

    def mousePressEvent(self, event):
        if self.mode == "none":
            super().mousePressEvent(event)
            return
            
        if event.button() == Qt.LeftButton:
            pos = event.position()
            self.current_path = QPainterPath()
            self.current_path.moveTo(pos)
            
            # Setup pen
            if self.mode == "pen":
                pen = QPen(QColor(255, 85, 85), 3, Qt.SolidLine, Qt.RoundCap, Qt.RoundJoin)
                self.paths.append({"path": self.current_path, "pen": pen})
            elif self.mode == "highlight":
                pen = QPen(QColor(255, 235, 59, 90), 14, Qt.SolidLine, Qt.RoundCap, Qt.RoundJoin)
                self.paths.append({"path": self.current_path, "pen": pen})
            elif self.mode == "eraser":
                self.erase_at(pos)
                
            self.update()

    def mouseMoveEvent(self, event):
        if self.mode == "none":
            super().mouseMoveEvent(event)
            return
            
        pos = event.position()
        if self.mode == "pointer":
            self.laser_pos = pos
            self.update()
            return
            
        if event.buttons() & Qt.LeftButton:
            if self.mode == "eraser":
                self.erase_at(pos)
            elif self.current_path:
                self.current_path.lineTo(pos)
                self.update()

    def mouseReleaseEvent(self, event):
        if self.mode == "none":
            super().mouseReleaseEvent(event)
            return
        self.current_path = None
        self.laser_pos = None
        self.update()

    def erase_at(self, pos):
        # Simple distance collision eraser
        new_paths = []
        for p in self.paths:
            collision = False
            # Check points along path
            for i in range(50):
                pt = p["path"].pointAtPercent(i / 50.0)
                if math.hypot(pt.x() - pos.x(), pt.y() - pos.y()) < 20:
                    collision = True
                    break
            if not collision:
                new_paths.append(p)
        self.paths = new_paths
        self.update()

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        
        # Draw lines
        for p in self.paths:
            painter.setPen(p["pen"])
            painter.drawPath(p["path"])
            
        # Draw laser pointer glow dot
        if self.mode == "pointer" and self.laser_pos:
            painter.setPen(Qt.NoPen)
            painter.setBrush(QBrush(QColor(255, 50, 50, 220)))
            painter.drawEllipse(self.laser_pos, 6, 6)
            painter.setBrush(QBrush(QColor(255, 50, 50, 60)))
            painter.drawEllipse(self.laser_pos, 14, 14)

    def clear_drawings(self):
        self.paths.clear()
        self.update()

    def resizeEvent(self, event):
        super().resizeEvent(event)


class PDFViewer(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.doc = None
        self.current_page = 0
        self.zoom = 1.2
        self.page_cache = {}
        
        self.setFocusPolicy(Qt.StrongFocus)
        self.init_ui()

    def init_ui(self):
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        
        # Toolbar
        toolbar = QHBoxLayout()
        toolbar.setContentsMargins(10, 10, 10, 10)
        
        self.btn_prev = QPushButton(" Previous")
        self.btn_next = QPushButton(" Next")
        self.lbl_page = QLabel("Page 0 of 0")
        self.lbl_page.setStyleSheet("color: #d4d4d4; font-weight: bold;")
        
        self.btn_zoom_in = QPushButton(" Zoom In")
        self.btn_zoom_out = QPushButton(" Zoom Out")
        
        self.btn_prev.setIcon(VectorIconProvider.get_icon("step_backward", "#ffffff"))
        self.btn_next.setIcon(VectorIconProvider.get_icon("step_forward", "#ffffff"))
        self.btn_zoom_in.setIcon(VectorIconProvider.get_icon("view_3d", "#007acc"))
        self.btn_zoom_out.setIcon(VectorIconProvider.get_icon("view_2d", "#858585"))
        
        toolbar.addWidget(self.btn_prev)
        toolbar.addWidget(self.lbl_page)
        toolbar.addWidget(self.btn_next)
        toolbar.addSpacing(20)
        
        # Overlay Drawing Toolbar for teaching
        self.btn_mouse = QPushButton(" Mouse")
        self.btn_pen = QPushButton(" Pen")
        self.btn_highlighter = QPushButton(" Highlight")
        self.btn_pointer = QPushButton(" Laser")
        self.btn_clear = QPushButton(" Clear")
        
        self.btn_mouse.setIcon(VectorIconProvider.get_icon("view_2d", "#4ec9b0"))
        self.btn_pen.setIcon(VectorIconProvider.get_icon("pen", "#ff5555"))
        self.btn_highlighter.setIcon(VectorIconProvider.get_icon("highlighter", "#ffb86c"))
        self.btn_pointer.setIcon(VectorIconProvider.get_icon("play", "#ff5555"))
        self.btn_clear.setIcon(VectorIconProvider.get_icon("clear", "#ff5555"))
        
        for btn in (self.btn_prev, self.btn_next, self.btn_zoom_in, self.btn_zoom_out,
                    self.btn_mouse, self.btn_pen, self.btn_highlighter, self.btn_pointer, self.btn_clear):
            btn.setStyleSheet("""
                QPushButton {
                    background-color: #333333;
                    color: #ffffff;
                    border: 1px solid #555555;
                    border-radius: 4px;
                    padding: 6px 12px;
                }
                QPushButton:hover {
                    background-color: #444444;
                }
                QPushButton:pressed {
                    background-color: #007acc;
                }
            """)
            
        self.btn_prev.clicked.connect(self.prev_page)
        self.btn_next.clicked.connect(self.next_page)
        self.btn_zoom_in.clicked.connect(self.zoom_in)
        self.btn_zoom_out.clicked.connect(self.zoom_out)
        
        self.btn_mouse.clicked.connect(lambda: self.set_overlay_mode("none"))
        self.btn_pen.clicked.connect(lambda: self.set_overlay_mode("pen"))
        self.btn_highlighter.clicked.connect(lambda: self.set_overlay_mode("highlight"))
        self.btn_pointer.clicked.connect(lambda: self.set_overlay_mode("pointer"))
        self.btn_clear.clicked.connect(self.clear_overlay_drawings)
        
        toolbar.addWidget(self.btn_mouse)
        toolbar.addWidget(self.btn_pen)
        toolbar.addWidget(self.btn_highlighter)
        toolbar.addWidget(self.btn_pointer)
        toolbar.addWidget(self.btn_clear)
        
        toolbar.addStretch()
        toolbar.addWidget(self.btn_zoom_out)
        toolbar.addWidget(self.btn_zoom_in)
        
        layout.addLayout(toolbar)
        
        # Display Area
        self.scroll_area = QScrollArea()
        self.scroll_area.setAlignment(Qt.AlignCenter)
        self.scroll_area.setStyleSheet("background-color: #252526; border: none;")
        
        self.lbl_image = QLabel()
        self.lbl_image.setAlignment(Qt.AlignCenter)
        self.scroll_area.setWidget(self.lbl_image)
        self.scroll_area.setWidgetResizable(True)
        
        layout.addWidget(self.scroll_area)
        
        # Setup Overlay Drawing Child Widget
        self.overlay = AnnotationOverlay(self.lbl_image)

    def set_overlay_mode(self, mode):
        self.overlay.set_mode(mode)
        # Highlight active tool button color
        for btn, m in ((self.btn_mouse, "none"), (self.btn_pen, "pen"), 
                       (self.btn_highlighter, "highlight"), (self.btn_pointer, "pointer")):
            if m == mode:
                btn.setStyleSheet("background-color: #007acc; color: white; border-radius: 4px; padding: 6px 12px;")
            else:
                btn.setStyleSheet("background-color: #333333; color: white; border-radius: 4px; padding: 6px 12px;")

    def clear_overlay_drawings(self):
        self.overlay.clear_drawings()

    def resizeEvent(self, event):
        super().resizeEvent(event)
        self.update_overlay_geometry()

    def update_overlay_geometry(self):
        self.overlay.setGeometry(self.lbl_image.rect())

    def load_pdf(self, path):
        if not fitz:
            self.lbl_image.setText("PyMuPDF is not installed. PDF viewing disabled.")
            self.lbl_image.setStyleSheet("color: red; font-size: 14px;")
            return
            
        try:
            self.doc = fitz.open(path)
            self.current_page = 0
            self.page_cache.clear()
            self.render_page()
        except Exception as e:
            self.lbl_image.setText(f"Error loading PDF: {str(e)}")
            self.lbl_image.setStyleSheet("color: red; font-size: 14px;")

    def render_page(self):
        if not self.doc:
            return
            
        if self.current_page in self.page_cache:
            pixmap = self.page_cache[self.current_page]
        else:
            page = self.doc.load_page(self.current_page)
            mat = fitz.Matrix(self.zoom, self.zoom)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            qimg = QImage.fromData(pix.tobytes("png"))
            pixmap = QPixmap.fromImage(qimg)
            self.page_cache[self.current_page] = pixmap
            
        self.lbl_image.setPixmap(pixmap)
        self.lbl_image.setFixedSize(pixmap.size())
        
        self.lbl_page.setText(f"Page {self.current_page + 1} of {len(self.doc)}")
        self.btn_prev.setEnabled(self.current_page > 0)
        self.btn_next.setEnabled(self.current_page < len(self.doc) - 1)
        
        self.overlay.clear_drawings()
        self.update_overlay_geometry()

    def prev_page(self):
        if self.current_page > 0:
            self.current_page -= 1
            self.render_page()

    def next_page(self):
        if self.doc and self.current_page < len(self.doc) - 1:
            self.current_page += 1
            self.render_page()

    def zoom_in(self):
        self.zoom += 0.15
        self.page_cache.clear()
        self.render_page()

    def zoom_out(self):
        if self.zoom > 0.4:
            self.zoom -= 0.15
            self.page_cache.clear()
            self.render_page()

    def keyPressEvent(self, event):
        # Keyboard navigation inside PDF
        if event.key() == Qt.Key_Left or event.key() == Qt.Key_PageUp:
            self.prev_page()
        elif event.key() == Qt.Key_Right or event.key() == Qt.Key_PageDown or event.key() == Qt.Key_Space:
            self.next_page()
        elif event.modifiers() == Qt.ControlModifier:
            if event.key() == Qt.Key_Equal or event.key() == Qt.Key_Plus:
                self.zoom_in()
            elif event.key() == Qt.Key_Minus:
                self.zoom_out()
        else:
            super().keyPressEvent(event)

    def wheelEvent(self, event):
        # Zoom with pinch gesture / Ctrl+Scroll
        if event.modifiers() == Qt.ControlModifier:
            if event.angleDelta().y() > 0:
                self.zoom_in()
            else:
                self.zoom_out()
        else:
            super().wheelEvent(event)


class PPTXViewer(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.slides = []
        self.current_slide = 0
        
        self.setFocusPolicy(Qt.StrongFocus)
        self.init_ui()

    def init_ui(self):
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        
        # Slide Display Widget
        self.slide_card = QWidget()
        self.slide_card.setObjectName("SlideCard")
        self.slide_card.setStyleSheet("""
            QWidget#SlideCard {
                background-color: #1e1e1e;
                border: 2px solid #3c3c3c;
                border-radius: 8px;
            }
        """)
        
        card_layout = QVBoxLayout(self.slide_card)
        card_layout.setContentsMargins(40, 40, 40, 40)
        
        self.lbl_title = QLabel("Slide Title")
        self.lbl_title.setWordWrap(True)
        self.lbl_title.setFont(QFont("Arial", 22, QFont.Bold))
        self.lbl_title.setStyleSheet("color: #007acc;")
        
        self.txt_content = QTextBrowser()
        self.txt_content.setStyleSheet("background-color: transparent; border: none; color: #d4d4d4;")
        self.txt_content.setFont(QFont("Arial", 14))
        
        card_layout.addWidget(self.lbl_title)
        card_layout.addWidget(self.txt_content)
        
        layout.addWidget(self.slide_card, 1)
        
        # Navigation Bar
        nav_layout = QHBoxLayout()
        nav_layout.setContentsMargins(10, 10, 10, 10)
        
        self.btn_prev = QPushButton(" Previous Slide")
        self.btn_next = QPushButton(" Next Slide")
        self.btn_prev.setIcon(VectorIconProvider.get_icon("step_backward", "#ffffff"))
        self.btn_next.setIcon(VectorIconProvider.get_icon("step_forward", "#ffffff"))
        
        self.lbl_slide_info = QLabel("Slide 0 of 0")
        self.lbl_slide_info.setStyleSheet("color: #d4d4d4; font-weight: bold;")
        
        # Presentation slide draw tools
        self.btn_mouse = QPushButton(" Mouse")
        self.btn_pen = QPushButton(" Pen")
        self.btn_highlighter = QPushButton(" Highlight")
        self.btn_pointer = QPushButton(" Laser")
        self.btn_clear = QPushButton(" Clear")
        
        self.btn_mouse.setIcon(VectorIconProvider.get_icon("view_2d", "#4ec9b0"))
        self.btn_pen.setIcon(VectorIconProvider.get_icon("pen", "#ff5555"))
        self.btn_highlighter.setIcon(VectorIconProvider.get_icon("highlighter", "#ffb86c"))
        self.btn_pointer.setIcon(VectorIconProvider.get_icon("play", "#ff5555"))
        self.btn_clear.setIcon(VectorIconProvider.get_icon("clear", "#ff5555"))
        
        for btn in (self.btn_prev, self.btn_next, self.btn_mouse, self.btn_pen, 
                    self.btn_highlighter, self.btn_pointer, self.btn_clear):
            btn.setStyleSheet("""
                QPushButton {
                    background-color: #333333;
                    color: #ffffff;
                    border: 1px solid #555555;
                    border-radius: 4px;
                    padding: 6px 12px;
                }
                QPushButton:hover {
                    background-color: #444444;
                }
                QPushButton:pressed {
                    background-color: #007acc;
                }
            """)
            
        self.btn_prev.clicked.connect(self.prev_slide)
        self.btn_next.clicked.connect(self.next_slide)
        
        self.btn_mouse.clicked.connect(lambda: self.set_overlay_mode("none"))
        self.btn_pen.clicked.connect(lambda: self.set_overlay_mode("pen"))
        self.btn_highlighter.clicked.connect(lambda: self.set_overlay_mode("highlight"))
        self.btn_pointer.clicked.connect(lambda: self.set_overlay_mode("pointer"))
        self.btn_clear.clicked.connect(self.clear_overlay_drawings)
        
        nav_layout.addWidget(self.btn_prev)
        nav_layout.addSpacing(25)
        nav_layout.addWidget(self.btn_mouse)
        nav_layout.addWidget(self.btn_pen)
        nav_layout.addWidget(self.btn_highlighter)
        nav_layout.addWidget(self.btn_pointer)
        nav_layout.addWidget(self.btn_clear)
        nav_layout.addStretch()
        nav_layout.addWidget(self.lbl_slide_info)
        nav_layout.addStretch()
        nav_layout.addWidget(self.btn_next)
        
        layout.addLayout(nav_layout)
        
        # Transparent annotation overlay on slides card
        self.overlay = AnnotationOverlay(self.slide_card)

    def set_overlay_mode(self, mode):
        self.overlay.set_mode(mode)
        for btn, m in ((self.btn_mouse, "none"), (self.btn_pen, "pen"), 
                       (self.btn_highlighter, "highlight"), (self.btn_pointer, "pointer")):
            if m == mode:
                btn.setStyleSheet("background-color: #007acc; color: white; border-radius: 4px; padding: 6px 12px;")
            else:
                btn.setStyleSheet("background-color: #333333; color: white; border-radius: 4px; padding: 6px 12px;")

    def clear_overlay_drawings(self):
        self.overlay.clear_drawings()

    def resizeEvent(self, event):
        super().resizeEvent(event)
        self.overlay.setGeometry(self.slide_card.rect())

    def load_pptx(self, path):
        if not Presentation:
            self.lbl_title.setText("python-pptx is not installed.")
            self.txt_content.setText("Please run: pip install python-pptx")
            return
            
        try:
            prs = Presentation(path)
            self.slides = []
            
            for s in prs.slides:
                title_text = ""
                body_paragraphs = []
                
                if s.shapes.title:
                    title_text = s.shapes.title.text
                
                for shape in s.shapes:
                    if shape.has_text_frame and shape != s.shapes.title:
                        for p in shape.text_frame.paragraphs:
                            t = p.text.strip()
                            if t:
                                body_paragraphs.append(t)
                                
                self.slides.append({
                    "title": title_text or "[Untitled Slide]",
                    "content": body_paragraphs
                })
                
            self.current_slide = 0
            self.render_slide()
        except Exception as e:
            self.lbl_title.setText("Failed to load Presentation")
            self.txt_content.setText(str(e))

    def render_slide(self):
        if not self.slides:
            return
            
        slide = self.slides[self.current_slide]
        self.lbl_title.setText(slide["title"])
        
        html = "<ul>"
        for bullet in slide["content"]:
            html += f"<li style='margin-bottom: 8px;'>{bullet}</li>"
        html += "</ul>"
        self.txt_content.setHtml(html)
        
        self.lbl_slide_info.setText(f"Slide {self.current_slide + 1} of {len(self.slides)}")
        self.btn_prev.setEnabled(self.current_slide > 0)
        self.btn_next.setEnabled(self.current_slide < len(self.slides) - 1)
        
        self.overlay.clear_drawings()
        self.overlay.setGeometry(self.slide_card.rect())

    def prev_slide(self):
        if self.current_slide > 0:
            self.current_slide -= 1
            self.render_slide()

    def next_slide(self):
        if self.slides and self.current_slide < len(self.slides) - 1:
            self.current_slide += 1
            self.render_slide()

    def keyPressEvent(self, event):
        if event.key() == Qt.Key_Left or event.key() == Qt.Key_PageUp:
            self.prev_slide()
        elif event.key() == Qt.Key_Right or event.key() == Qt.Key_PageDown or event.key() == Qt.Key_Space:
            self.next_slide()
        else:
            super().keyPressEvent(event)


class DOCXViewer(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.init_ui()

    def init_ui(self):
        layout = QVBoxLayout(self)
        layout.setContentsMargins(10, 10, 10, 10)
        
        self.browser = QTextBrowser()
        self.browser.setStyleSheet("""
            QTextBrowser {
                background-color: #1e1e1e;
                color: #d4d4d4;
                border: 1px solid #3c3c3c;
                border-radius: 4px;
                padding: 20px;
            }
        """)
        self.browser.setFont(QFont("Segoe UI", 12))
        layout.addWidget(self.browser)

    def load_docx(self, path):
        if not Document:
            self.browser.setHtml("<h2 style='color: red;'>python-docx is not installed.</h2>")
            return
            
        try:
            doc = Document(path)
            html = []
            
            html.append("<style>")
            html.append("h1 { color: #007acc; border-bottom: 1px solid #3c3c3c; padding-bottom: 5px; }")
            html.append("h2 { color: #4ec9b0; }")
            html.append("h3 { color: #ce9178; }")
            html.append("p { line-height: 1.5; margin-bottom: 12px; }")
            html.append("</style>")
            
            for p in doc.paragraphs:
                text = p.text.strip()
                if not text:
                    continue
                    
                style_name = p.style.name if p.style else ""
                
                if style_name.startswith("Heading"):
                    try:
                        level = int(style_name[-1])
                    except:
                        level = 1
                    html.append(f"<h{level}>{text}</h{level}>")
                else:
                    para_html = ""
                    for run in p.runs:
                        run_text = run.text
                        if run.bold:
                            run_text = f"<b>{run_text}</b>"
                        if run.italic:
                            run_text = f"<i>{run_text}</i>"
                        para_html += run_text
                    html.append(f"<p>{para_html}</p>")
                    
            self.browser.setHtml("\n".join(html))
        except Exception as e:
            self.browser.setHtml(f"<h2 style='color: red;'>Failed to load Document</h2><p>{str(e)}</p>")


class PresentationTab(QWidget):
    """
    Unified widget that loads any document type.
    Handles global keyboard listeners for page flips.
    """
    def __init__(self, parent=None):
        super().__init__(parent)
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        
        self.stack = QStackedWidget()
        layout.addWidget(self.stack)
        
        self.pdf_viewer = PDFViewer()
        self.pptx_viewer = PPTXViewer()
        self.docx_viewer = DOCXViewer()
        
        self.stack.addWidget(self.pdf_viewer)
        self.stack.addWidget(self.pptx_viewer)
        self.stack.addWidget(self.docx_viewer)
        
        self.fallback = QLabel("Select a document from the resource explorer to view.")
        self.fallback.setAlignment(Qt.AlignCenter)
        self.fallback.setStyleSheet("color: #858585; font-size: 14px;")
        self.stack.addWidget(self.fallback)
        self.stack.setCurrentIndex(3)

    def load_file(self, path):
        ext = os.path.splitext(path)[1].lower()
        if ext == ".pdf":
            self.stack.setCurrentIndex(0)
            self.pdf_viewer.load_pdf(path)
            self.pdf_viewer.setFocus() # Focus keyboard events
        elif ext in (".pptx", ".ppt"):
            self.stack.setCurrentIndex(1)
            self.pptx_viewer.load_pptx(path)
            self.pptx_viewer.setFocus()
        elif ext in (".docx", ".doc"):
            self.stack.setCurrentIndex(2)
            self.docx_viewer.load_docx(path)
        else:
            self.stack.setCurrentIndex(3)
            self.fallback.setText(f"Unsupported file format: {ext}")
